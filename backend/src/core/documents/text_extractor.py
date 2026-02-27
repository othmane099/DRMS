import logging
from pathlib import Path

logger = logging.getLogger(__name__)

_SUPPORTED = {".pdf", ".docx", ".xlsx", ".pptx", ".txt"}
_SKIPPED = {".doc", ".xls", ".jpg", ".jpeg", ".png"}


async def extract_text(file_path: str) -> str | None:
    """Return extracted text from a document file, or None if unsupported/empty."""
    path = Path(file_path)
    suffix = path.suffix.lower()

    if suffix in _SKIPPED:
        return None

    if suffix not in _SUPPORTED:
        return None

    try:
        text = _extract(path, suffix)
    except Exception:
        logger.exception("Text extraction failed for %s", file_path)
        return None

    if not text:
        return None

    return text


def _extract(path: Path, suffix: str) -> str | None:
    if suffix == ".txt":
        return path.read_text(errors="replace")

    if suffix == ".pdf":
        import pypdf

        reader = pypdf.PdfReader(str(path))
        parts = []
        for page in reader.pages:
            try:
                parts.append(page.extract_text() or "")
            except Exception:
                pass
        return "\n".join(parts)

    if suffix == ".docx":
        import docx

        doc = docx.Document(str(path))
        return "\n".join(p.text for p in doc.paragraphs)

    if suffix == ".xlsx":
        import openpyxl

        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        parts = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_text = " ".join(str(cell) for cell in row if cell is not None)
                if row_text:
                    parts.append(row_text)
        return "\n".join(parts)

    if suffix == ".pptx":
        from pptx import Presentation

        prs = Presentation(str(path))
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            parts.append(text)
        return "\n".join(parts)

    return None
